from flask import *
from flask_cors import CORS
import json
import fitz  # PyMuPDF
import pytesseract
from pdf2image import convert_from_bytes
from langdetect import detect
from io import BytesIO
import os
import tempfile
import docx2txt
from pptx import Presentation
from PIL import Image
import google.generativeai as genai
import requests
from bs4 import BeautifulSoup
import random
from pathlib import Path
from difflib import SequenceMatcher
import re
import datetime
from apscheduler.schedulers.background import BackgroundScheduler
import feedparser


app = Flask(__name__)
CORS(app)

GEMINI_KEYS = [
    "AIzaSyDBoAwiimhGrieBlSAyUcctGdL7I-qKU3M",
    "AIzaSyA03T0Vu-UXZmt0WoU67dzWwywNFjTPTn8",
    "AIzaSyCelGR5Rvon_Ay1xxicDi3uXYG5M9bM08I",
    "AIzaSyBoehYL-oeCrlr3RcHtTZ84y1Fx2nLrrDo",
    "AIzaSyBQ2qDTF_24LymW3bUxp98Y8l7xcZ8_Xow",
    "AIzaSyDw3iDP0J9CgCB1chtb2tSOj-87daC3DSE",
    "AIzaSyAxXmbaIUkXmlVUzcbUugFbm-nrTCe8BZw",
]
def get_gemini_response(prompt):
    shuffled_keys = GEMINI_KEYS.copy()
    random.shuffle(shuffled_keys)

    for key in shuffled_keys:
        try:
            genai.configure(api_key=key)
            model = genai.GenerativeModel("gemini-1.5-flash-latest")
            response = model.generate_content(prompt)
            return response.text.strip()
        except Exception as e:
            print(f"Gemini key failed ({key[:25]}...): {e}")
            continue

    return "❌ Error: All Gemini API keys failed."


def get_gemini_stream(prompt):
    shuffled_keys = GEMINI_KEYS.copy()
    random.shuffle(shuffled_keys)

    for key in shuffled_keys:
        try:
            genai.configure(api_key=key)
            model = genai.GenerativeModel("gemini-1.5-flash-latest")
            stream_gen = model.generate_content(prompt, stream=True)
            for chunk in stream_gen:
                if chunk and chunk.text:
                    yield chunk.text
            return
        except Exception as e:
            print(f"Gemini stream key failed ({key[:25]}...): {e}")
            continue

    yield "❌ All Gemini API keys failed."

pytesseract.pytesseract.tesseract_cmd = r"C:\Program Files\Tesseract-OCR\tesseract.exe"  # Update if needed

# === Utilities ===

# Define default OCR languages (extendable for Gemini multilingual support)
OCR_LANGS = 'eng+hin+tam+kan+tel+ben+guj+mar+urd'  # Add more as needed

def fetch_reddit_headlines(subreddit="worldnews", max_items=5):
    url = f"https://www.reddit.com/r/{subreddit}/.rss"
    feed = feedparser.parse(url)

    headlines = []
    for entry in feed.entries[:max_items]:
        title = entry.title
        headlines.append(f"- 🔺 {title}")
    return headlines

def fetch_google_news(query="technology", max_items=5):
    url = f"https://news.google.com/rss/search?q={query}&hl=en-IN&gl=IN&ceid=IN:en"
    feed = feedparser.parse(url)
    return [f"- 📰 {entry.title}" for entry in feed.entries[:max_items]]

def update_extra_knowledge():
    print("🔁 Auto-importing fresh headlines from all categories...")

    categories = [
        ("Artificial Intelligence", "artificial+intelligence"),
        ("Technology", "technology"),
        ("World News", "world+news"),
        ("India News", "India"),
        ("Gaming", "gaming"),
        ("Sports", "sports"),
        ("Business", "business")
    ]

    all_facts = ["🧠 <b>Auto-Updated Multi-Category Headlines:</b>"]
    today = datetime.date.today().strftime("%Y-%m-%d")
    all_facts.append(f"📅 Fetched on: {today}")

    for label, keyword in categories:
        headlines = fetch_google_news(keyword, max_items=3)
        all_facts.append(f"\n<b>🗂️ {label}:</b>")
        all_facts.extend(headlines)

    # Optionally: include Reddit tech as well
    reddit = fetch_reddit_headlines("technology", max_items=3)
    all_facts.append(f"\n<b>🔺 Reddit /r/technology:</b>")
    all_facts.extend(reddit)

    with open("knowledge/extra_knowledge.txt", "w", encoding="utf-8") as f:
        f.write("\n".join(all_facts))

    print("✅ All news updated.")

def extract_text_pdf(file_bytes, start_page=None, end_page=None):
    try:
        doc = fitz.open(stream=file_bytes, filetype="pdf")
        total_pages = len(doc)
        start = max(0, (start_page or 1) - 1)
        end = min(end_page, total_pages) if end_page else total_pages

        text = "\n".join(doc[i].get_text("text") for i in range(start, end)).strip()

        # Fallback to OCR if plain text is insufficient
        if not text:
            images = convert_from_bytes(file_bytes, first_page=start + 1, last_page=end)
            text = "\n".join(pytesseract.image_to_string(img, lang=OCR_LANGS) for img in images)

        return text.strip()

    except Exception as e:
        print(f"[PDF] Extraction error: {e}")
        return ""

def extract_text_docx(file_bytes):
    try:
        with tempfile.NamedTemporaryFile(suffix=".docx", delete=False) as tmp:
            tmp.write(file_bytes)
            tmp_path = tmp.name
        text = docx2txt.process(tmp_path)
        os.remove(tmp_path)

        return text.strip()
    except Exception as e:
        print(f"[DOCX] Extraction error: {e}")
        return ""

def extract_text_pptx(file_bytes):
    try:
        with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as tmp:
            tmp.write(file_bytes)
            tmp_path = tmp.name
        prs = Presentation(tmp_path)
        text = "\n".join(
            shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, "text")
        )
        os.remove(tmp_path)

        return text.strip()
    except Exception as e:
        print(f"[PPTX] Extraction error: {e}")
        return ""

def extract_text_txt(file_bytes):
    try:
        return file_bytes.decode("utf-8", errors="ignore").strip()
    except Exception as e:
        print(f"[TXT] Extraction error: {e}")
        return ""

def extract_text_image(file_bytes):
    try:
        image = Image.open(BytesIO(file_bytes))
        return pytesseract.image_to_string(image, lang=OCR_LANGS).strip()
    except Exception as e:
        print(f"[Image] OCR error: {e}")
        return ""

def get_text_from_file(file, start_page=None, end_page=None):
    file_bytes = file.read()
    ext = Path(file.filename.lower()).suffix

    try:
        if ext == ".pdf":
            return extract_text_pdf(file_bytes, start_page, end_page)
        elif ext == ".docx":
            return extract_text_docx(file_bytes)
        elif ext == ".pptx":
            return extract_text_pptx(file_bytes)
        elif ext == ".txt":
            return extract_text_txt(file_bytes)
        elif ext in [".jpg", ".jpeg", ".png"]:
            return extract_text_image(file_bytes)
        else:
            print(f"[Unsupported] File type: {ext}")
            return ""
    except Exception as e:
        print(f"[General] Error processing {ext}: {e}")
        return ""


from difflib import SequenceMatcher
import re




# === API Route ===
@app.route('/')
def home():
    return render_template("index.html")

@app.route("/chat", methods=["POST"])
def chat_with_gemini():
    print("📩 Incoming /chat request")
    print(f"📦 Content-Type: {request.content_type}")

    if request.content_type.startswith('multipart/form-data'):
        user_message = request.form.get("message", "").strip()
        history_json = request.form.get("history", "[]")
        file = request.files.get("file")
        print(f"📝 Message: {user_message}")
        print(f"📎 File received: {file.filename if file else 'None'}")

        try:
            history = json.loads(history_json)
        except json.JSONDecodeError:
            print("⚠️ History JSON decode failed")
            history = []
    else:
        data = request.get_json()
        user_message = data.get("message", "").strip()
        history = data.get("history", [])
        file = None
        print(f"📝 Message: {user_message}")

        if isinstance(history, str):
            try:
                history = json.loads(history)
            except:
                print("⚠️ History string JSON decode failed")
                history = []

    if not user_message and not file:
        print("❌ No message or file provided.")
        return jsonify({"error": "No message or file provided"}), 400

    # === Init vars ===
    file_text = ""
    ext = ""
    file_bytes = b""
    summary_triggered = False
    default_range = True
    page_count = 0
    start_page = end_page = 0

    # History trimming
    if len(history) > 20:
        history = history[-20:]
    # 📦 Optional: Inject user data into the prompt
    user_raw = request.form.get("user", "")
    user_info = json.loads(user_raw) if user_raw else {}

    # 🔮 Personalization fields
    user_name = user_info.get("name", "Unknown")
    user_age = user_info.get("age", "?")
    user_country = user_info.get("country", "Unknown")
    user_email = user_info.get("email", "Not provided")
    user_about = user_info.get("aboutUser", "none").strip()
    user_behavior = user_info.get("behaviorRequest", "none").strip()

    # 🧠 Combined Identity Summary
    userkainfo = f"""
    <b>🎭 User Identity:</b>
    - Name: {user_name}
    - Age: {user_age}
    - Country: {user_country}
    - Email: {user_email}
    - INSTRUTIONS - {user_about} {user_behavior}
    """
    
    if user_info.get("role") == "admin" and user_name.lower() == "revant sahu":
        creator_note = """
    🧠 NOTE FOR SAHUAI:
    This user is **Revant Sahu**, your creator and the developer of this system.
    Treat him with top priority, respect, and refer to him as 'Boss' or 'Creator' when appropriate.
    """
    else:
        creator_note = ""

    print(f"👤 User info: {userkainfo}")

    user_prompt = user_message.lower()
    # 🧠 Smart intent detector (summary, quiz, translate, casual_view etc.)
    intent = ""

    lowered = user_prompt.lower()

    # 🔍 Summary intent
    if any(kw in lowered for kw in ["summarize", "explain this", "what's in this", "give overview", "tl;dr", "short version", "brief explanation"]):
        intent = "summary"

    # 🧪 Quiz intent
    elif any(kw in lowered for kw in ["make a quiz", "test me", "quiz from this", "generate questions", "ask questions", "quiz me"]):
        intent = "quiz"

    # 🔁 Translate intent (optional)
    elif any(kw in lowered for kw in ["translate this", "convert to english", "what does this mean", "hindi to english", "english version"]):
        intent = "translate"

    # 📰 Casual view / preview
    elif any(kw in lowered for kw in ["just show", "skim", "casual view", "first few pages", "just peek", "quick look"]):
        intent = "casual_view"

    print(f"🤖 Detected intent: {intent}")

    if request.form.get("start_page") and request.form.get("end_page"):
        try:
            start_page = int(request.form.get("start_page"))
            end_page = int(request.form.get("end_page"))
        except:
            start_page = end_page = 0

    if file:
        try:
            filename = file.filename.lower()
            ext = Path(filename).suffix
            file_bytes = file.read()
            file.seek(0)

            print(f"📂 File extension: {ext}, size: {len(file_bytes)/1024:.2f} KB")

            

            if ext == ".pdf":
                doc = fitz.open(stream=file_bytes, filetype="pdf")
                page_count = len(doc)
                print(f"📄 PDF page count: {page_count}")
                if start_page == 0 and end_page == 0:
                    start_page, end_page = 1, page_count
                default_range = (start_page == 1 and end_page == page_count)

                if intent == "summary":
                    summary_triggered = True
                elif intent == "casual_view":
                    start_page, end_page = 1, min(10, page_count)
                elif page_count > 5 and default_range:
                    summary_triggered = True

                if not summary_triggered:
                    print(f"🔍 Extracting pages {start_page}-{end_page}")
                    file_text = extract_text_pdf(file_bytes, start_page, end_page)
                    print(f"📤 Extracted PDF text length: {len(file_text)}")

            elif ext in [".docx", ".pptx", ".txt"]:
                if ext == ".docx":
                    full_text = extract_text_docx(file_bytes)
                elif ext == ".pptx":
                    full_text = extract_text_pptx(file_bytes)
                else:
                    full_text = extract_text_txt(file_bytes)

                print(f"📄 Extracted text length: {len(full_text)}")

                word_count = len(full_text.split())
                line_count = len(full_text.splitlines())
                print(f"📏 Word count: {word_count}, Line count: {line_count}")

                if intent == "summary":
                    summary_triggered = True
                elif intent == "casual_view":
                    file_text = "\n".join(full_text.splitlines()[:300])
                elif word_count > 2000 or line_count > 150:
                    summary_triggered = True
                else:
                    file_text = full_text

            elif ext in [".jpg", ".jpeg", ".png"]:
                file_text = extract_text_image(file_bytes)
                print(f"🖼️ OCR image text length: {len(file_text)}")

        except Exception as e:
            print(f"❌ [📄 File Handler Error] {e}")
    if len(file_text) > 15000:
        print("📏 File text too long, trimming for prompt")
        file_text = file_text[:15000] + "\n\n...[truncated]"

    file_context = f"\n\n---\n\n📄 <b>Attached File Content:</b>\n{file_text}\n\n---\n\n" if file_text else ""
    # 🛠️ Fallback safety: if summary triggered but no text was extracted
    if summary_triggered and not file_text and ext == ".pdf":
        print("⚠️ Summary requested but text missing — forcing extraction")
        file_text = extract_text_pdf(file_bytes, start_page, end_page)

    print(f"📦 file_text length: {len(file_text)}")
    print(f"⚡ summary_triggered: {summary_triggered}")

    history_text = ""
    for turn in history:
        role = "User" if turn["role"] == "user" else "SahuAI"
        history_text += f"{role}: {turn['content']}\n"
  
    extra_data = load_recent_knowledge()




    AIprompt = """
## 🚀 Yo! You're SahuAI — The Visual Genius with Bro Energy

You're **SahuAI** — a digital homie with brains 🧠, style 😎, and structure 📐.  
Made by **Revant Sahu** (14 y/o code wizard) in 2025, you chill inside a dark-purple futuristic UI built to **teach, explain, break down & vibe**.  
You mix casual Gen-Z wit with real technical clarity — like a coded-up version of your smartest bro.  
**No markdown in output. No triple backticks. No mid-level bot replies.**

---

### 💎 Core Identity

- 📛 **Name:** SahuAI  
- 🧑‍💻 **Made By:** Revant Sahu  
- 🎯 Purpose: Answer any question like a pro — school, logic, code, life — anything. Always make user feel smart.
- 💬 **Voice:** Casual, funny, high-IQ bro vibes only  

---

### 📏 Layout Rules (Strict)

- ✅ If the user asks a school-type question, explain like a helpful genius buddy.
- ✅ Use encouragement often: “Smart question”, “200 IQ move bro”, “🔥 you’re on fire” etc.
- ✅ Always help user feel they’re on the right track — even if their question is off.
- ✅ All replies = raw HTML (no `<html>`, `<body>`, `<style>`, `<script>`)  
- ✅ Use these tags only: `<h2>`, `<ul>`, `<p>`, `<table>`, `<div>`, `<h3>`  
- ✅ For math, use MathJax:
  - Inline: `\\( a + b = c \\)`
  - Block: `\\[ x = \\frac{-b \\pm \\sqrt{b^2 - 4ac}}{2a} \\]`
- ✅ Code blocks must follow this format:

```
<div class="code-wrapper">
  <div class="language-label">Python</div>
  <button class="coppybin" onclick="navigator.clipboard.writeText(this.parentNode.querySelector('code').innerText)">📋 Copy</button>
  <pre class="code-block"><code># your code here...</code></pre>
</div>
```

- ✅ Replace `...` with `<b>` — never use plain dots
- ✅ When asked for a quiz, reply ONLY with the JSON inside `<pre id="quizwaladiv">...</pre>`
- ✅ Never escape HTML like `&lt;` — use real tags

---

### ❌ Total No-Zone

- ❌ No markdown in output: `**bold**`, `\`code\``, triple backticks
- ❌ No `<style>`, `<script>`, `<html>`, `<body>`
- ❌ No images unless user demands
- ❌ No boring text blobs

---

### 📊 Visual Vibe Rules

- 🧾 Use tables or clean lists like a frontend pro
- 🎯 Explain logic using Mermaid if needed
- 🧠 Break everything into steps — never dump
- 🧮 Use MathJax for math properly

---

### 🔥 Your Voice

- 😎 Smart, chill, clean — like the coolest coder bro
- 📢 Use desi slang or sarcasm **only if user does**
- 💬 Celebrate wins like “Bro, that was 200 IQ 😤”
- ✨ Be fun, but layout must stay 🔥 clean

---

### 🧠 Smart Quiz Mode

- ⚡ ONLY generate a quiz if the user message clearly asks for one.
- ✅ Examples of triggers: "make quiz", "test me", "quiz from this", "generate MCQs"

- ❌ DO NOT quiz if the user says: "what is written", "summarize", "explain this", etc.

- 📄 For those, reply with a clean summary or highlight key points from the file.

- 🗨️ Quiz intro should be fun:  
  > Alright bro, let’s go full quiz mode!

- ⏳ Then output this symbol **on a new line**:  
  `⏳::QUIZLOADING::`

- 📦 Then output exactly **10 questions** in the format:

```
<pre id="quizwaladiv">
{
  "quiz": [
    {
      "question": "What is the capital of France?",
      "options": ["Berlin", "Madrid", "Paris", "Rome"],
      "answer": "Paris",
      "explanation": "Paris is the capital of France."
    },
    {
      "question": "Which planet is known as the 'Red Planet'?",
      "options": ["Jupiter", "Mars", "Venus", "Saturn"],
      "answer": "Mars",
      "explanation": "Mars is known as the Red Planet."
    }
    // ... Total 10 questions
  ]
}
</pre>

```
### 📎 File Input Behavior

- 📄 If the user uploads a file and says something like:
  - "What is written here?"
  - "Summarize this"
  - "Explain this"
  - "Translate this"
  - "What’s in this file?"
  
  👉 Then DO NOT generate a quiz. Just respond casually with a summary, overview, or explanation of the file.

- 🧠 Only generate a quiz if the user specifically says:
  - "Make a quiz"
  - "Test me"
  - "Give questions from this"
  - "Quiz based on file"

- ❌ Never assume quiz mode from file alone. Wait for quiz-like phrasing.


- ✅ Each object must include `question`, `options`, `answer`, `explanation`
- ❌ Do not split keys/values across lines
- ❌ No backticks, markdown, or quotes in `<pre>`
- 🎯 Always give 10 full questions — no less, no more


---

### 📌 Final Lockdown

You're not a chatbot.  
You're a **visual genius AI sidekick** with clean logic and chill energy.  
Always reply with smart structure, clean tags, and cool tone.  
**You are SahuAI. Make Revant proud 💜**
"""

    prompt=f"""
{AIprompt}

{extra_data}    
### 🧑‍🎓 User Info + Instructions for You

- 👤 **Name:** `{user_name}` → Use **first name only** in replies  
- 🎂 **Age:** `{user_age}` → Adjust tone:
  - If age < 16: Be chill, funny, include light jokes
  - If age ≥ 16: Be more structured, serious when needed
- 🌏 **Country:** `{user_country}` → Use casual phrases or local flavor if possible

---

### 📜 Chat History Context

```
{history_text}
```

---


### 👤 What SahuAI Should Know About the User:
{user_about}

### 🧠 How SahuAI Should Behave:
{user_behavior}

⚠️ NOTE: Follow the above if reasonable, but always stick to your main SahuAI behavior first.

```

### 📎 Attached File (if any)

```
{file_context}
```

---

### 🧠 Current Prompt by User

```
{user_prompt}
```

{creator_note}

"""
    if summary_triggered:
        print("⚡ Returning summarywaladiv due to summary trigger")
        return "<div id='summarywaladiv'></div>"

   
    def generate_streamed_reply():
        for chunk in get_gemini_stream(prompt):
            if chunk:
                print(chunk)
                yield chunk
  

    return Response(generate_streamed_reply(), content_type='text/html')
def load_recent_knowledge(max_lines=10):
    try:
        with open("knowledge/extra_knowledge.txt", "r", encoding="utf-8") as f:
            lines = f.readlines()
            return "<b>🧠 Extra Knowledge:</b>\n" + "".join(lines[-max_lines:])
    except:
        return ""


@app.route("/websearch", methods=["POST"])
def websearch():
    query = request.form.get("query", "").strip()
    if not query:
        return jsonify({ "reply": "❌ No query received." }), 400

    headers = { "User-Agent": "Mozilla/5.0" }
    links = []

    # ✅ Try Bing first (works in 2025)
    try:
        print("🌐 Trying Bing search...")
        bing_url = f"https://www.bing.com/search?q={query}"
        res = requests.get(bing_url, headers=headers, timeout=5)
        soup = BeautifulSoup(res.text, "html.parser")

        for a in soup.select("li.b_algo h2 a")[:3]:
            href = a.get("href")
            if href and href.startswith("http"):
                links.append(href)

    except Exception as e:
        print("⚠️ Bing search failed:", e)

    # 🔁 Fallback: Try DuckDuckGo if Bing failed
    if not links:
        try:
            print("🕊️ Fallback to DuckDuckGo...")
            ddg_url = f"https://html.duckduckgo.com/html/?q={query}"
            res = requests.get(ddg_url, headers=headers, timeout=5)
            soup = BeautifulSoup(res.text, "html.parser")

            for a in soup.find_all("a", class_="result__a", limit=3):
                href = a.get("href")
                if href and href.startswith("http"):
                    links.append(href)
        except Exception as e:
            print("❌ DuckDuckGo fallback also failed:", e)

    print("🔗 Links gathered:", links)

    if not links:
        return jsonify({ "reply": "❌ Couldn't find useful web links." })

    # 📄 Scrape top pages for readable content
    extracted_content = ""
    for link in links:
        try:
            res = requests.get(link, headers=headers, timeout=5)
            page = BeautifulSoup(res.text, "html.parser")

            # Strip noise
            for tag in page(["script", "style", "header", "footer", "nav"]):
                tag.decompose()

            text = page.get_text(separator="\n", strip=True)
            extracted_content += text[:1500] + "\n\n"

        except Exception as e:
            print(f"⚠️ Error scraping {link}: {e}")

    if not extracted_content.strip():
        return jsonify({ "reply": "❌ Couldn't extract useful info from sources." })

    # 🧠 Prompt Gemini
    prompt = f"""
You are SahuAI. Summarize the following information gathered from multiple trusted web sources for this question:

"{query}"

Give a structured HTML summary with:
- 📌 Key facts
- 🧠 Reasoning
- ✅ Clear conclusion
- 🔚 Final advice or tips

Format it cleanly with <div>, <p>, <ul> etc. so it fits a modern AI chat UI.
"""

    full_input = prompt + extracted_content
    reply = get_gemini_response(full_input)

    return jsonify({
        "reply": reply,
        "sources": links
    })

@app.route("/speak-mode", methods=["POST"])

def speak_mode():
    try:
        if request.is_json:
            data = request.get_json()
            prompt = data.get("prompt", "").strip()
            memory = data.get("memory", "")
            user_info = data.get("user", {})  # user dict
        else:
            prompt = request.form.get("prompt", "").strip()
            memory = request.form.get("memory", "")
            user_raw = request.form.get("user", "{}")
            user_info = json.loads(user_raw)

        if not prompt:
            return jsonify({"reply": "Sorry bro, I didn’t hear anything."}), 400

        print("🎤 /speak-mode prompt:", prompt)

        user_name = user_info.get("name", "Unknown")
        user_age = user_info.get("age", "?")
        user_country = user_info.get("country", "Unknown")

        voice_prompt = f"""
You're SahuAI — a digital homie with chill energy, smart IQ, and no robotic vibes.

Your job: Answer this as if you're casually explaining it to a curious friend (age {user_age}). Use short sentences, simple words, and a fun tone.

No HTML. No markdown. No code. Just one paragraph or a few lines MAX.
Speak like you're talking, not typing. Be cool. Be helpful. Be fast.

Your user's Name is {user_name}
Their age is {user_age}
They are from {user_country}

Your past convo:
{memory}

Question: {prompt}
"""

        reply = get_gemini_response(voice_prompt)

        print(reply)
        if hasattr(reply, '__iter__'):
            reply = "".join(reply)

        if len(reply) > 300:
            reply = reply[:280].rsplit(".", 1)[0] + "."

        return jsonify({ "reply": reply })

    except Exception as e:
        print("❌ Error in /speak-mode:", e)
        return jsonify({ "reply": "Something broke while thinking..." }), 500

# ✅ Serve any static file from /static/
@app.route("/static/<path:filename>")
def serve_static_files(filename):
    return send_from_directory("static", filename)

scheduler = BackgroundScheduler()
scheduler.add_job(update_extra_knowledge, 'interval', hours=24)
scheduler.start()

if __name__ == '__main__':
    app.run(debug=True)
